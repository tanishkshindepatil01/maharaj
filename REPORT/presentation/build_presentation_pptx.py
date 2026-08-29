import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # University of Bristol Theme Colors
    C_RED = RGBColor(176, 28, 46)        # Bristol Red (#B01C2E)
    C_DARK_RED = RGBColor(140, 20, 36)   # Deep Burgundy (#8C1424)
    C_NAVY = RGBColor(12, 35, 64)        # Deep Navy (#0C2340)
    C_CHARCOAL = RGBColor(30, 41, 59)    # Dark Slate Charcoal (#1E293B)
    C_MUTED = RGBColor(100, 116, 139)    # Slate Gray (#64748B)
    C_BG = RGBColor(248, 250, 252)       # Light Off-White (#F8FAFC)
    C_WHITE = RGBColor(255, 255, 255)
    C_CARD_BG = RGBColor(255, 255, 255)
    C_BORDER = RGBColor(226, 232, 240)   # Light border (#E2E8F0)
    C_LIGHT_RED = RGBColor(253, 242, 242)

    logo_path = 'logo_uob_color.png' if os.path.exists('logo_uob_color.png') else 'REPORT/latex_source 2/logo_uob_color.png'

    def add_header(slide, slide_num, category, title):
        # Top banner line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
        line.fill.solid()
        line.fill.fore_color.rgb = C_RED
        line.line.color.rgb = C_RED

        # Category / Section tracker + Slide number
        tb_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(9.0), Inches(0.3))
        tf_cat = tb_cat.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = f"{slide_num:02d}  |  {category.upper()}"
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = C_RED

        # Title
        tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(10.0), Inches(0.65))
        tf_title = tb_title.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = C_NAVY

        # Small UoB text or logo in corner
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(10.8), Inches(0.25), width=Inches(1.8))

        # Bottom subtle footer
        tb_foot = slide.shapes.add_textbox(Inches(0.8), Inches(7.1), Inches(11.7), Inches(0.3))
        tf_foot = tb_foot.text_frame
        p_foot = tf_foot.paragraphs[0]
        p_foot.text = "University of Bristol  •  MSc Data Science Dissertation Project  •  McCain Foods & Dalhousie Collaboration"
        p_foot.font.size = Pt(9)
        p_foot.font.color.rgb = C_MUTED

    def add_card(slide, left, top, width, height, bg_color=C_CARD_BG, border_color=C_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
        return card

    # ==========================================
    # SLIDE 1: TITLE SLIDE (MATCHING WHITE & RED UNIVERSITY COLORS)
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    # White / Light Off-White Background
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = C_BG
    bg1.line.fill.background()

    # Top red accent bar
    top_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = C_RED
    top_bar.line.fill.background()

    # Left vertical red accent stripe
    accent1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(0.1), Inches(2.2))
    accent1.fill.solid()
    accent1.fill.fore_color.rgb = C_RED
    accent1.line.fill.background()

    # University of Bristol Logo on Top Left
    if os.path.exists(logo_path):
        s1.shapes.add_picture(logo_path, Inches(0.8), Inches(0.6), width=Inches(3.2))

    # Title box
    tb = s1.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(11.2), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Edge-Deployable Rock Detection\n& Burial-Status Classification"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = C_NAVY

    p2 = tf.add_paragraph()
    p2.text = "Real-Time Embedded Computer Vision for Precision Agriculture on Raspberry Pi 5"
    p2.font.size = Pt(18)
    p2.font.color.rgb = C_RED
    p2.space_before = Pt(8)

    # Info card with white background and subtle border
    card_authors = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.3))
    card_authors.fill.solid()
    card_authors.fill.fore_color.rgb = C_WHITE
    card_authors.line.color.rgb = C_BORDER
    card_authors.line.width = Pt(1.2)

    tb_meta = s1.shapes.add_textbox(Inches(1.1), Inches(4.6), Inches(11.1), Inches(2.1))
    tf_meta = tb_meta.text_frame
    tf_meta.word_wrap = True

    # STRICT AUTHOR HIERARCHY
    pm1 = tf_meta.paragraphs[0]
    pm1.text = "Project Team: 1. Tanishk Shinde  •  2. Megh Kashilkar  •  3. Sanskar Gade  •  4. Jitendra Suwalka"
    pm1.font.size = Pt(14)
    pm1.font.bold = True
    pm1.font.color.rgb = C_NAVY

    pm2 = tf_meta.add_paragraph()
    pm2.text = "Academic Supervisor: Dr. Felipe Campelo (University of Bristol)"
    pm2.font.size = Pt(12)
    pm2.font.color.rgb = C_CHARCOAL
    pm2.space_before = Pt(6)

    pm3 = tf_meta.add_paragraph()
    pm3.text = "External Supervisors: Dr. Ahmad Al-Mallahi & Jason McLaren (Dalhousie University / McCain Foods)"
    pm3.font.size = Pt(12)
    pm3.font.color.rgb = C_CHARCOAL

    pm4 = tf_meta.add_paragraph()
    pm4.text = "University of Bristol  •  Faculty of Engineering  •  Department of Engineering Mathematics  •  16 September 2026"
    pm4.font.size = Pt(11)
    pm4.font.color.rgb = C_MUTED
    pm4.space_before = Pt(6)

    s1.notes_slide.notes_text_frame.text = "Welcome everyone. This presentation details our MSc Data Science project on Edge-Deployable Rock Detection and Burial Status Classification, presented by Tanishk Shinde, Megh Kashilkar, Sanskar Gade, and Jitendra Suwalka."

    # Helper function for 2-column slides with cards
    def create_two_column_slide(slide_num, category, title, left_card_data, right_card_data, note="", font_size=13, space_before=8):
        slide = prs.slides.add_slide(blank_layout)
        add_header(slide, slide_num, category, title)

        # Left Card
        add_card(slide, 0.8, 1.3, 5.7, 5.5)
        tb_l = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(5.3), Inches(5.1))
        tf_l = tb_l.text_frame
        tf_l.word_wrap = True
        
        pl0 = tf_l.paragraphs[0]
        pl0.text = left_card_data['title']
        pl0.font.size = Pt(18)
        pl0.font.bold = True
        pl0.font.color.rgb = C_RED

        for pt in left_card_data['points']:
            p = tf_l.add_paragraph()
            p.text = "• " + pt
            p.font.size = Pt(font_size)
            p.font.color.rgb = C_CHARCOAL
            p.space_before = Pt(space_before)

        # Right Card
        add_card(slide, 6.8, 1.3, 5.7, 5.5)
        if right_card_data.get('image') and os.path.exists(right_card_data['image']):
            tb_r = slide.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.3), Inches(0.5))
            tf_r = tb_r.text_frame
            tf_r.word_wrap = True
            pr0 = tf_r.paragraphs[0]
            pr0.text = right_card_data['title']
            pr0.font.size = Pt(18)
            pr0.font.bold = True
            pr0.font.color.rgb = C_NAVY

            slide.shapes.add_picture(right_card_data['image'], Inches(7.0), Inches(2.1), width=Inches(5.3))
            if right_card_data.get('caption'):
                tb_c = slide.shapes.add_textbox(Inches(7.0), Inches(6.1), Inches(5.3), Inches(0.6))
                tf_c = tb_c.text_frame
                tf_c.word_wrap = True
                pc = tf_c.paragraphs[0]
                pc.text = right_card_data['caption']
                pc.font.size = Pt(10)
                pc.font.color.rgb = C_MUTED
        else:
            tb_r = slide.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.3), Inches(5.1))
            tf_r = tb_r.text_frame
            tf_r.word_wrap = True
            pr0 = tf_r.paragraphs[0]
            pr0.text = right_card_data['title']
            pr0.font.size = Pt(18)
            pr0.font.bold = True
            pr0.font.color.rgb = C_NAVY

            for pt in right_card_data.get('points', []):
                p = tf_r.add_paragraph()
                p.text = "• " + pt
                p.font.size = Pt(font_size)
                p.font.color.rgb = C_CHARCOAL
                p.space_before = Pt(space_before)

        if note:
            slide.notes_slide.notes_text_frame.text = note
        return slide

    # SLIDE 2: Operational Problem
    create_two_column_slide(
        2, "Context & Motivation", "Rocks Cause Severe Operational & Equipment Costs in Potato Farming",
        {
            'title': "Agricultural Need & Operational Impact",
            'points': [
                "Commercial potato harvesting machinery suffers severe mechanical wear and blade damage from unmapped field rocks.",
                "Manual rock picking is labour-intensive, hazardous, and causes costly harvesting bottlenecks during peak season.",
                "Precision removal requires detecting not just rock location, but also burial status (completely exposed vs. half-buried) to apply appropriate mechanical force.",
                "Dalhousie University and McCain Foods initiated this project to enable autonomous on-the-go vision systems directly on field machinery."
            ]
        },
        {
            'title': "Industry Collaboration Objective",
            'points': [
                "Partner: McCain Foods & Dalhousie University Agricultural Engineering.",
                "Goal: Deploy real-time object detection directly on tractor-mounted edge hardware.",
                "Constraint: Real field environments lack cloud connectivity, demanding fully offline edge inference.",
                "Challenge: Variable soil appearance, partial occlusions, and severe compute/power limitations on embedded devices."
            ]
        },
        "Rocks create severe mechanical damage and downtime in commercial potato farming. Our goal is to enable real-time detection on tractor-mounted hardware."
    )

    # SLIDE 3: Edge-Deployment Challenge
    create_two_column_slide(
        3, "Context & Motivation", "High Model Accuracy Does Not Guarantee Real-Time Edge Feasibility",
        {
            'title': "The Academic vs. Edge Deployment Gap",
            'points': [
                "Standard vision research prioritises heavy GPU models (e.g. YOLOv8x, ViT) and high precision (FP32).",
                "On-the-go tractor operations operate at ~5–10 km/h, requiring at least 15–20 FPS to guide hydraulic actuators in real time.",
                "PyTorch models on low-cost single-board computers (Raspberry Pi 5) achieve only 2.5 FPS (~400 ms latency) — completely unacceptable for real-time control.",
                "Quantisation (FP16, INT8) and edge-optimised runtimes (TFLite, NCNN, ONNX) are mandatory, but their accuracy–speed–stability trade-offs were previously unbenchmarked."
            ]
        },
        {
            'title': "Key Technical Constraints",
            'points': [
                "Compute: Low-power quad-core ARM Cortex-A76 processor without discrete GPU acceleration.",
                "Memory: Peak RAM must remain well within device limits without paging or out-of-memory crashes.",
                "Thermal: Continuous field operation without active cooling must prevent thermal CPU throttling.",
                "Reliability: The runtime must maintain 100% stability across continuous hours without segmentation faults."
            ]
        },
        "Standard high-accuracy models fail on edge devices due to latency. We need an end-to-end edge deployment pipeline."
    )

    # SLIDE 4: Research Question & Hypotheses
    create_two_column_slide(
        4, "Context & Motivation", "Core Research Question & Empirical Hypotheses",
        {
            'title': "The Central Research Question",
            'points': [
                "\"How effectively can lightweight YOLO models be optimised, quantised, and deployed for real-time rock detection and burial-status classification on a low-cost edge device (Raspberry Pi 5)?\"",
                "Sub-question 1: How does model architecture selection (YOLOv8n vs YOLO11n vs YOLO11s) impact detection of challenging half-buried rocks?",
                "Sub-question 2: What is the exact empirical trade-off between numerical precision (FP32 vs FP16 vs INT8) in terms of latency, accuracy (mAP50-95), and system stability?"
            ]
        },
        {
            'title': "Guiding Empirical Hypotheses",
            'points': [
                "H1 (Speedup): Quantised edge-runtime representations (TFLite INT8 / NCNN) will yield at least a 5x latency reduction compared to PyTorch FP32 baseline.",
                "H2 (Accuracy Cost): INT8 quantisation will incur a measurable localisation degradation (mAP50-95 drop < 5%) due to 8-bit dynamic range compression.",
                "H3 (Operational Robustness): Deployment formats will exhibit non-trivial divergence in CPU, thermal, and stability behaviours under continuous operation."
            ]
        },
        "We formulated a precise research question and three testable hypotheses around speedup, accuracy cost, and operational stability."
    )

    # SLIDE 5: Start-to-Finish Pipeline
    create_two_column_slide(
        5, "Methodology Overview", "Rigorous Five-Stage End-to-End Project Architecture",
        {
            'title': "Project Execution Workflow",
            'points': [
                "1. Data Integrity & Splitting: Pre-augmentation split into 313 train, 45 val, and 90 test images to prevent data leakage.",
                "2. Controlled Augmentation: 6 safe photometric transforms expanding the training set to 2,191 images (3,563 rock instances).",
                "3. Controlled Model Selection: Identical training of YOLOv8n, YOLO11n, and YOLO11s; selection strictly based on validation mAP50-95.",
                "4. Six-Format Edge Export: Exporting frozen YOLOv8n into PyTorch, TorchScript, ONNX, TFLite (FP32/INT8), and NCNN (FP32/FP16).",
                "5. Systematic Pi 5 Benchmark: Fair, letterboxed 90-image test set evaluation measuring FPS, latency, mAP, count MAE, RAM, CPU, and thermals."
            ]
        },
        {
            'title': "Controlled Variables & Principles",
            'points': [
                "Zero Test-Set Leakage: Test set (90 images) never augmented and never used during model selection.",
                "Frozen Weights: Identical neural network weights exported across all six runtime formats.",
                "Identical Letterboxing: Standardized 640x640 input with symmetric fill value (114) for all backends.",
                "Repeatable Protocol: 3 warmup runs followed by timed inference with hardware telemetry logging."
            ]
        },
        "Our methodology follows five strict stages ensuring no data leakage, fair model selection, and standardized edge evaluation.",
        font_size=10.5,
        space_before=3
    )

    # SLIDE 6: Dataset Characteristics
    create_two_column_slide(
        6, "Dataset Preparation", "Original Field Dataset: 448 Images & 706 Annotated Rocks",
        {
            'title': "Dataset Composition & Class Balance",
            'points': [
                "448 high-resolution aerial and ground-level field images collected from real potato cultivation plots.",
                "706 manually verified bounding box annotations across two critical operational classes:",
                "• Completely Exposed: 362 objects (51.3%) — clear surface boundary, high visual contrast against soil.",
                "• Half-Buried: 344 objects (48.7%) — partially occluded by soil clods, irregular silhouettes, lower contrast.",
                "Near 50:50 class balance allows performance differences to be interpreted as visual difficulty rather than class imbalance."
            ]
        },
        {
            'title': "Dataset Distribution Summary",
            'image': 'figures/figure_04.png',
            'caption': "Figure 3.4: Dataset summary showing 448 original images, 706 annotated rocks, and 2 balanced burial classes."
        },
        "The dataset has 448 images and 706 rocks, cleanly balanced between completely exposed and half-buried classes."
    )

    # SLIDE 7: Splitting First Protects Evidence
    create_two_column_slide(
        7, "Dataset Preparation", "Pre-Augmentation Splitting Prevents Data Leakage",
        {
            'title': "Partitioning Strategy",
            'points': [
                "Partitioning was performed BEFORE any augmentation occurred to ensure zero duplicate variants leaked into evaluation.",
                "• Training Partition: 313 images (259 exposed, 250 half-buried) — 69.9% of dataset.",
                "• Validation Partition: 45 images (34 exposed, 28 half-buried) — 10.0% used exclusively for architecture selection.",
                "• Test Partition: 90 images (69 exposed, 66 half-buried) — 20.1% untouched held-out edge benchmark set.",
                "Preserves absolute independence for all subsequent Raspberry Pi runtime comparisons."
            ]
        },
        {
            'title': "Class Splits Across Partitions",
            'image': 'figures/figure_02.png',
            'caption': "Figure 3.2: Class-level annotation counts across Train, Val, and Test splits showing balanced proportions."
        },
        "We partitioned the data prior to augmentation to guarantee that no test or validation image had augmented counterparts in training."
    )

    # SLIDE 8: Safe Augmentation
    create_two_column_slide(
        8, "Dataset Preparation", "Controlled Photometric Augmentation Expands Training to 2,191 Images",
        {
            'title': "Safe Offline Augmentation Strategy",
            'points': [
                "Each of the 313 training images generated 6 photometric variants (+1,878 images), expanding training set to 2,191 images.",
                "Total training annotation instances increased from 509 to 3,563 rock objects.",
                "Augmentations simulated real field illumination shifts:",
                "• Brightness & Contrast adjustment (cloud cover vs bright direct sunlight)",
                "• Gamma correction & HSV colour jitter (dry vs moist soil tones)",
                "• Gaussian noise & slight blur (camera motion / vibration)",
                "Geometric distortion (shear, extreme crop) was strictly avoided to prevent invalidating bounding box boundaries."
            ]
        },
        {
            'title': "Training Set Augmentation Expansion",
            'image': 'figures/figure_03.png',
            'caption': "Figure 3.3: 7x expansion of training data from 313 original images to 2,191 augmented images."
        },
        "We expanded training data safely using 6 photometric transforms per image, avoiding geometric distortions that corrupt bounding boxes."
    )

    # SLIDE 9: Single-Stage YOLO Architecture
    create_two_column_slide(
        9, "Model Architecture", "Single-Stage YOLO Detects Location & Burial Status Simultaneously",
        {
            'title': "Why Single-Stage YOLO for Field Edge AI?",
            'points': [
                "Two-stage detectors (Faster R-CNN) require separate region proposals, resulting in prohibitively high latency (>800 ms on ARM).",
                "Single-stage YOLO architectures predict bounding boxes [x, y, w, h] and class probabilities in a single forward pass.",
                "Unified loss function optimizes bounding box regression (CIoU / DFL) alongside classification loss (BCE).",
                "Anchor-free design in modern YOLO (v8, 11) enhances small rock detection and adapts to irregular rock shapes without manual anchor tuning."
            ]
        },
        {
            'title': "Edge Output Tensor Representation",
            'points': [
                "Input: Standardized 640x640x3 RGB image tensor.",
                "Output: 8,400 candidate anchor predictions decoded into [x, y, w, h, p_exposed, p_buried].",
                "Post-Processing:",
                "• Confidence Threshold: Filtering candidates with confidence < 0.25.",
                "• Non-Maximum Suppression (NMS): IoU threshold 0.70 to eliminate duplicate detections.",
                "• Counting Aggregator: Outputs total detected rocks, exposed count, and half-buried count per field frame."
            ]
        },
        "YOLO performs simultaneous localization and classification in a single pass, ideal for low-latency edge deployment."
    )

    # SLIDE 10: Controlled Candidate Training
    create_two_column_slide(
        10, "Model Development", "Controlled Training of Lightweight YOLO Candidates",
        {
            'title': "Rigorous Experimental Controls",
            'points': [
                "Three candidate lightweight architectures evaluated: YOLOv8n (3.2M params), YOLO11n (2.6M params), and YOLO11s (9.4M params).",
                "Common training setup across all models to ensure fair comparison:",
                "• Epochs: 20 epochs with early checkpoint saving",
                "• Image Resolution: 896x896 training, 1024x1024 validation",
                "• Optimizer: AdamW with Cosine learning rate scheduling",
                "• Batch Size: 32 with Automatic Mixed Precision (AMP)",
                "• Pre-trained Weights: COCO transfer learning initialization"
            ]
        },
        {
            'title': "Precision, Recall & F1 Scores",
            'image': 'figures/figure_05.png',
            'caption': "Figure 3.5: Validation & test precision, recall, and F1 comparison across YOLOv8n, YOLO11n, and YOLO11s."
        },
        "All candidate models were trained under strictly identical hyperparameters to prevent tuning bias from favoring any single model."
    )

    # SLIDE 11: Architecture Selection Criterion
    create_two_column_slide(
        11, "Model Selection", "Predefined Primary Criterion: Validation mAP50–95",
        {
            'title': "Avoiding Post-Hoc Metric Cherry-Picking",
            'points': [
                "Architecture selection was governed by a predefined primary metric: Validation mAP50–95.",
                "Why mAP50–95 over mAP50?",
                "• mAP50 averages IoU threshold at 0.50 only — loose localisation boundary.",
                "• mAP50–95 averages IoU from 0.50 to 0.95 in 0.05 steps, strictly punishing sloppy box localisation.",
                "For precision agriculture, accurate boundary estimation is crucial for mechanical rock pickers to avoid striking rock edges.",
                "The test set was strictly locked during this phase."
            ]
        },
        {
            'title': "Candidate Validation Performance",
            'points': [
                "Validation Results (45 validation images):",
                "• YOLOv8n: Precision 0.681, Recall 0.759, F1 0.718, mAP50 0.759, mAP50-95 = 0.4288 (WINNER)",
                "• YOLO11n: Precision 0.743, Recall 0.719, F1 0.731, mAP50 0.729, mAP50-95 = 0.4002",
                "• YOLO11s: Precision 0.764, Recall 0.801, F1 0.782, mAP50 0.783, mAP50-95 = 0.4107",
                "Decision: YOLOv8n selected as the definitive detector for all edge deployment experiments."
            ]
        },
        "YOLOv8n won the predefined validation mAP50-95 criterion with 0.4288, making it our frozen checkpoint for edge export."
    )

    # SLIDE 12: Detector Selection Result
    create_two_column_slide(
        12, "Model Selection", "YOLOv8n Achieves Highest Validation Localisation Quality",
        {
            'title': "Validation Evaluation Summary",
            'points': [
                "YOLOv8n achieved mAP50–95 of 0.4288, outperforming YOLO11n (0.4002) by +7.1% and YOLO11s (0.4107) by +4.4%.",
                "Although YOLO11s scored slightly higher on F1 (0.782 vs 0.718) and mAP50 (0.783 vs 0.759), its 9.4M parameter size tripled compute requirements.",
                "YOLOv8n offered the optimal balance: top localisation precision with minimal 3.2M parameter footprint.",
                "Checkpoint weights were frozen and passed to the multi-format edge export pipeline."
            ]
        },
        {
            'title': "mAP50 & mAP50–95 Comparison",
            'image': 'figures/figure_06.png',
            'caption': "Figure 3.6: mAP50 and mAP50–95 comparison across candidate YOLO models on validation and test partitions."
        },
        "YOLOv8n delivered superior localisation precision while maintaining a compact 3.2M parameter footprint."
    )

    # SLIDE 13: Generalisation & Burial Difficulty
    create_two_column_slide(
        13, "Model Diagnostics", "Class Breakdown Reveals the Inherent Difficulty of Half-Buried Rocks",
        {
            'title': "Class-Specific Performance Asymmetry",
            'points': [
                "Evaluating YOLOv8n on the held-out 90-image test set revealed a significant disparity between classes:",
                "• Completely Exposed: 64 / 69 rocks detected correctly (92.8% operational accuracy).",
                "• Half-Buried: Only 23 / 66 rocks detected correctly (34.8% operational accuracy).",
                "Visual Causes of Half-Buried Failures:",
                "• Smooth gradual transitions between rock surfaces and surrounding dried soil.",
                "• Partial occlusion by mud clods, crop residue, and shadow contours.",
                "• Absence of distinct continuous bounding edges.",
                "Proves that burial-status classification is the primary computer vision bottleneck in field settings."
            ]
        },
        {
            'title': "Key Takeaway for Field Operations",
            'points': [
                "Exposed rocks are easily detected with >90% reliability, sufficient to trigger fast surface-sweep mechanisms.",
                "Half-buried rocks require deeper soil penetration and higher actuation forces; misclassification leads to mechanical jams.",
                "Future multi-spectral or depth sensing will be required to elevate half-buried detection rates to industrial levels."
            ]
        },
        "Exposed rocks achieved 92.8% detection, whereas half-buried rocks achieved only 34.8% due to visual soil occlusion."
    )

    # SLIDE 14: Auxiliary CLIP Investigation
    create_two_column_slide(
        14, "Auxiliary Investigation", "CLIP Enhances Crop Classification but Lacks Native Localization",
        {
            'title': "Zero-Shot vs Fine-Tuned Vision-Language Model",
            'points': [
                "Investigated whether foundation models (OpenAI CLIP ViT-B/32) could enhance burial status discrimination.",
                "Experimental Setup: Cropped rock image chips evaluated under zero-shot prompting vs fine-tuned linear probe.",
                "Results on Rock Chips:",
                "• Zero-Shot CLIP: 50.4% accuracy (essentially random guess between exposed vs buried).",
                "• Fine-Tuned CLIP: 80.0% accuracy on cropped chips.",
                "Why CLIP Was Not Deployed on Edge:",
                "• CLIP is a whole-image/crop classifier, NOT a native object detector; requires pre-extracted bounding boxes.",
                "• ViT-B/32 requires ~350 MB RAM and ~850 ms inference per crop on Raspberry Pi 5 — prohibitively slow."
            ]
        },
        {
            'title': "Role in the Project Scope",
            'points': [
                "Confirmed that foundation vision-language models have strong semantic feature representations once fine-tuned.",
                "However, single-stage YOLO remains vastly superior for joint localisation, classification, and edge throughput.",
                "CLIP findings documented as a valuable auxiliary diagnostic in Chapter 3."
            ]
        },
        "Fine-tuned CLIP reached 80% accuracy on cropped chips, but was not viable for edge detection due to heavy compute demands."
    )

    # SLIDE 15: Edge Inference Pipeline
    create_two_column_slide(
        15, "Edge Pipeline", "End-to-End Edge Pipeline: Letterboxing to Count Telemetry",
        {
            'title': "Optimized Edge Data Flow",
            'points': [
                "1. Input Normalisation: 640x640 letterboxing with aspect ratio preservation and symmetric fill (value 114).",
                "2. Tensor Pre-processing: RGB normalization [0, 1] and NCHW format conversion.",
                "3. Runtime Engine Inference: Multi-backend forward execution (PyTorch, TFLite, NCNN, ONNX).",
                "4. Anchor Decoding: Extracting box coordinates and multi-class confidence scores.",
                "5. Post-Processing & NMS: Threshold filtering (conf >= 0.25, IoU = 0.70).",
                "6. Agricultural Metric Telemetry: Frame latency, FPS, memory usage, total rock count, and per-class counts."
            ]
        },
        {
            'title': "Counting Accuracy Metric",
            'points': [
                "Mean Absolute Error (MAE) on total rock counts per image:",
                "• PyTorch FP32: MAE = 0.600 rocks/image",
                "• TFLite FP32: MAE = 0.600 rocks/image",
                "• TFLite INT8: MAE = 0.633 rocks/image (+0.033 rock count delta)",
                "• NCNN FP32: MAE = 0.600 rocks/image",
                "• NCNN FP16: MAE = 0.611 rocks/image",
                "Shows that INT8 quantization maintains excellent aggregate field rock counting fidelity."
            ]
        },
        "The edge pipeline letterboxes images to 640x640, runs inference, applies NMS, and outputs bounding boxes, labels, and count MAE."
    )

    # SLIDE 16: Qualitative Error Analysis
    create_two_column_slide(
        16, "Model Diagnostics", "Qualitative Failure Modes: Soil Clods & Occlusion Boundaries",
        {
            'title': "Systematic Error Pattern Analysis",
            'points': [
                "False Positives (Soil Clods):",
                "• Dry, crusted soil clumps frequently exhibit convex shapes and shadow gradients mimicking small rocks.",
                "• Can be mitigated by elevating confidence threshold (e.g. 0.35) at the cost of slight recall reduction.",
                "False Negatives (Deeply Buried Rocks):",
                "• Rocks with <20% surface exposure blend into surrounding tillage patterns.",
                "Class Confusion:",
                "• Edge boundary uncertainty: Rocks with ~50% burial often oscillate between exposed and half-buried classes.",
                "Conclusion: 2D RGB imagery alone reaches an empirical performance ceiling; 3D surface profiling is recommended for future work."
            ]
        },
        {
            'title': "Visual Diagnostic Lessons",
            'points': [
                "Illumination Invariance: Photometric augmentation successfully prevented false alarms under harsh direct sunlight.",
                "Scale Sensitivity: Rocks < 30 pixels across have lower detection confidence; 640x640 resolution represents optimal speed-resolution trade-off.",
                "All visual failure cases thoroughly categorized to guide industrial implementation."
            ]
        },
        "Failure analysis shows soil clods cause false positives and deep burial causes false negatives, indicating 2D RGB limits."
    )

    # SLIDE 17: Edge Experiment Preparation
    create_two_column_slide(
        17, "Edge Methodology", "Preparation of Frozen Checkpoint & Representative Calibration Set",
        {
            'title': "Standardized Edge Preparation Protocol",
            'points': [
                "The selected YOLOv8n weights were permanently frozen into a master PyTorch checkpoint (`best.pt`, 6.2 MB).",
                "90 Test Images Standardized: Letterboxed to 640x640 format, creating a fixed evaluation testbed.",
                "INT8 Calibration Dataset Preparation:",
                "• Extracted 200 original training images containing 304 rock objects.",
                "• Strictly balanced: exactly 152 completely exposed and 152 half-buried instances.",
                "• Used to compute static scale factor S and zero-point Z for integer-only quantisation without backpropagation.",
                "Ensured zero test data was observed during quantisation calibration."
            ]
        },
        {
            'title': "Affine INT8 Quantisation Formula",
            'points': [
                "Quantisation mapping: q = round(r / S) + Z",
                "• r: 32-bit floating point activation/weight",
                "• q: 8-bit signed integer [-128, 127]",
                "• S: positive real scaling factor",
                "• Z: integer zero-point offset",
                "Enables ARM NEON integer dot-product SIMD instructions, unlocking massive throughput gains."
            ]
        },
        "We prepared a frozen checkpoint, 90 standardized test images, and a balanced 200-image calibration set for INT8 quantization."
    )

    # SLIDE 18: Six Exported Representations
    create_two_column_slide(
        18, "Edge Methodology", "Six Deployment Representations Span FP32, FP16 & INT8",
        {
            'title': "The Six Evaluated Edge Representations",
            'points': [
                "1. PyTorch FP32 (6.2 MB): Native baseline running in PyTorch Python runtime.",
                "2. TorchScript FP32 (12.3 MB): JIT-compiled PyTorch intermediate representation.",
                "3. ONNX FP32 (12.2 MB): Open Neural Network Exchange format with ONNX Runtime.",
                "4. TFLite FP32 (12.0 MB): TensorFlow Lite flatbuffer single-precision candidate.",
                "5. TFLite INT8 (3.3 MB): TensorFlow Lite 8-bit post-training integer quantized model (-73% file size).",
                "6. NCNN FP32 (12.1 MB) & FP16 (6.1 MB): Tencent NCNN embedded framework for mobile ARM CPUs."
            ]
        },
        {
            'title': "Candidate Export Summary",
            'points': [
                "All models generated from identical YOLOv8n weights.",
                "Memory footprint drops from 12.3 MB (TorchScript) down to 3.3 MB (TFLite INT8).",
                "Enables direct empirical comparison across runtime engines (PyTorch vs TFLite vs NCNN vs ONNX) and numerical precisions (FP32 vs FP16 vs INT8)."
            ]
        },
        "We exported the same YOLOv8n detector into six deployment formats spanning PyTorch, ONNX, TFLite, and NCNN across FP32, FP16, and INT8."
    )

    # SLIDE 19: Raspberry Pi 5 Testbed
    create_two_column_slide(
        19, "Edge Methodology", "Hardware Testbed: Raspberry Pi 5 (8GB) in CPU-Only Mode",
        {
            'title': "Raspberry Pi 5 Hardware & OS Environment",
            'points': [
                "SoC: Broadcom BCM2712 quad-core ARM Cortex-A76 processor @ 2.4 GHz.",
                "Memory: 8 GB LPDDR4X SDRAM.",
                "Power Supply: Official 27W (5.1V / 5.0A) USB-C power delivery to prevent under-voltage throttling.",
                "Operating System: Debian GNU/Linux 13 (Trixie), 64-bit kernel.",
                "Cooling Condition: Passive heatsink only (No Active Cooler attached) to simulate realistic enclosed tractor equipment enclosures.",
                "Software Stack: Python 3.13.5, Ultralytics 8.4.115, ONNX Runtime 1.28.0, OpenCV 5.0.0, NumPy 2.5.1, psutil 7.2.2."
            ]
        },
        {
            'title': "Experimental Isolation Protocols",
            'points': [
                "No GUI Desktop: Background services minimized to eliminate operating system jitter.",
                "Single Batch Inference: Batch size = 1 strictly enforced (real-time stream simulation).",
                "Hardware Telemetry: psutil and vcgencmd polled every 10 ms for CPU load, RSS RAM, and SoC core temperature."
            ]
        },
        "Our testbed is a Raspberry Pi 5 with 8GB RAM running headless Debian Linux, actively monitoring CPU, memory, and thermal state."
    )

    # SLIDE 20: Benchmarking Protocol
    create_two_column_slide(
        20, "Edge Methodology", "Dual Testing Protocol: Fixed Benchmark & Continuous Endurance",
        {
            'title': "Two-Phase Evaluation Protocol",
            'points': [
                "Phase 1: Canonical Fixed Benchmark (90 Images)",
                "• 3 warmup inferences to stabilize cache and JIT compilation.",
                "• 90 timed sequential inferences over the held-out test partition.",
                "• Records per-frame latency, median FPS, CPU %, peak RAM, mAP50, mAP50-95, and Count MAE.",
                "Phase 2: Continuous Operational Endurance Test (10 Passes)",
                "• 900 consecutive inferences (10 passes of the 90-image test set).",
                "• Evaluates memory leak accumulation, thermal throttling thresholds, and runtime crash vulnerability."
            ]
        },
        {
            'title': "Why Dual Testing Was Essential",
            'points': [
                "A model can pass a short 90-image benchmark but crash after 20 minutes in field conditions.",
                "Exposes hidden edge runtime bugs (memory leaks, threading lockups, segmentation faults).",
                "Ensures the selected model is truly robust for continuous tractor deployments."
            ]
        },
        "We implemented both a 90-image canonical benchmark and a 900-image continuous endurance run to evaluate operational stability."
    )

    # SLIDE 21: Latency & Throughput Results (THE WINNER)
    create_two_column_slide(
        21, "Evaluation Results", "TFLite INT8 Delivers 22.62 FPS — 9.01x Speedup over PyTorch",
        {
            'title': "Canonical Runtime Latency & Throughput",
            'points': [
                "PyTorch FP32 (Baseline): 397.48 ms latency | 2.52 FPS — unusable for real-time control.",
                "TFLite FP32: 143.70 ms latency | 6.95 FPS (2.77x speedup over PyTorch).",
                "NCNN FP32: 75.09 ms latency | 13.26 FPS (5.29x speedup).",
                "NCNN FP16: 72.31 ms latency | 13.82 FPS (5.49x speedup).",
                "TFLite INT8 (WINNER): 44.13 ms latency | 22.62 FPS (9.01x speedup!).",
                "TFLite INT8 is the ONLY representation that exceeds the 20 FPS threshold required for real-time potato harvesting."
            ]
        },
        {
            'title': "Throughput & Latency Comparison",
            'image': 'figures/figure_09.png',
            'caption': "Figure 4.3: Median throughput (FPS) and mean latency (ms) across completed deployment formats on Pi 5."
        },
        "TFLite INT8 achieved 22.62 FPS and 44.13 ms latency — a 9.01x speedup over PyTorch and the only format exceeding the 20 FPS real-time target."
    )

    # SLIDE 22: Accuracy vs Throughput Trade-off
    create_two_column_slide(
        22, "Evaluation Results", "Throughput vs Accuracy: A Quantifiable, Acceptable Trade-Off",
        {
            'title': "Localisation Quality vs Inference Speed",
            'points': [
                "Full-Precision Quality Retention:",
                "• PyTorch FP32: mAP50 = 0.548, mAP50-95 = 0.275, F1 = 0.622",
                "• TFLite FP32: mAP50 = 0.535, mAP50-95 = 0.270, F1 = 0.599",
                "• NCNN FP16: mAP50 = 0.537, mAP50-95 = 0.270, F1 = 0.600",
                "INT8 Quantisation Impact (TFLite INT8):",
                "• mAP50: 0.503 (-0.045 vs PyTorch baseline)",
                "• mAP50-95: 0.234 (-0.041 vs PyTorch baseline)",
                "• F1 Score: 0.594 (-0.028 vs PyTorch baseline)",
                "Trade-Off Assessment: A 14.9% decrease in mAP50-95 buys an 800% increase in frame rate (2.5 -> 22.6 FPS). Highly favorable for field deployment."
            ]
        },
        {
            'title': "Speed vs Accuracy Pareto Chart",
            'image': 'figures/figure_deployment_speed_accuracy_scatter.png',
            'caption': "Figure 4.2: Median throughput (FPS) plotted against mAP50-95 showing TFLite INT8 as the clear Pareto-optimal choice."
        },
        "The minor mAP reduction from INT8 quantization is heavily outweighed by the massive 9x speedup, establishing TFLite INT8 as Pareto-optimal."
    )

    # SLIDE 23: Resource & Thermal Diagnostics
    create_two_column_slide(
        23, "Evaluation Results", "System Resources: INT8 Optimizes CPU, Memory & Thermal Load",
        {
            'title': "Hardware Resource Utilization on Pi 5",
            'points': [
                "Peak Process RAM:",
                "• TFLite INT8 had the lowest memory footprint: 824.2 MB (vs 868.3 MB for TFLite FP32).",
                "• Fits easily within low-cost 2GB or 4GB Raspberry Pi variants.",
                "CPU Core Utilization:",
                "• TFLite INT8: 62.8% average CPU load across 4 cores.",
                "• NCNN FP16: 87.1% CPU load (heavier ARM thread saturation).",
                "Thermal Dissipation (Short-run Max):",
                "• TFLite INT8: 55.10 °C (Coolest operating temperature).",
                "• NCNN FP32: 56.75 °C | NCNN FP16: 59.50 °C | PyTorch: 61.70 °C | TFLite FP32: 66.10 °C.",
                "INT8's integer SIMD operations generate significantly less heat, preventing thermal throttling."
            ]
        },
        {
            'title': "Resource & Thermal Comparison",
            'image': 'figures/figure_11.png',
            'caption': "Figure 4.5: CPU load, peak RAM, and short-run temperature across completed formats."
        },
        "TFLite INT8 was also the most power-efficient format, running coolest at 55.1°C with only 62.8% CPU load."
    )

    # SLIDE 24: Operational Stability & Negative Results
    create_two_column_slide(
        24, "Operational Stability", "Stability Testing Exposes Critical Runtime Failures",
        {
            'title': "Negative Results & Runtime Crashes",
            'points': [
                "Export Success Does NOT Equal Deployment Stability!",
                "ONNX FP32 Fatal Failure:",
                "• Successfully exported from Ultralytics without errors.",
                "• Completely crashed during benchmark aggregation due to an internal ONNX Runtime memory management fault on ARM64.",
                "NCNN Continuous Endurance Failure:",
                "• Successfully passed the 90-image benchmark at 13.8 FPS.",
                "• However, crashed during the continuous 10-pass endurance run due to thread contention and memory leakage."
            ]
        },
        {
            'title': "Robustness & Deployment Findings",
            'points': [
                "TFLite INT8 Robustness Winner:",
                "• The ONLY deployment format that completed 100% of both fixed benchmark and continuous endurance runs without a single crash.",
                "System-Level Validation Necessity:",
                "• Demonstrates the indispensable need to test models under sustained workload stress on target edge hardware.",
                "Reporting Integrity:",
                "• Transparently disclosing runtime crashes prevents costly real-world field deployment failures."
            ]
        },
        "Crucially, ONNX crashed during benchmarking and NCNN crashed during endurance testing. TFLite INT8 was the only format with 100% stability."
    )

    # SLIDE 25: Final Decision & Roadmap
    create_two_column_slide(
        25, "Conclusions & Roadmap", "Final Selection: TFLite INT8 Validated; Future Roadmap Defined",
        {
            'title': "Final Deployment Selection: TFLite INT8",
            'points': [
                "Selection Justification across All Four Criteria:",
                "1. Speed: 22.62 FPS (9.01x faster than PyTorch, exceeding 20 FPS threshold).",
                "2. Localisation: Retains 0.234 mAP50-95 and 0.633 count MAE.",
                "3. Efficiency: Lowest peak RAM (824 MB) and lowest temperature (55.1 °C).",
                "4. Reliability: 100% stability across all continuous endurance workloads.",
                "Hypothesis Assessment: All three hypotheses (H1 Speedup, H2 Accuracy Cost, H3 Stability Divergence) fully supported."
            ]
        },
        {
            'title': "Prioritised Future Roadmap",
            'points': [
                "1. Multi-Spectral & 3D Depth Sensing: Integrate stereo vision or LiDAR to elevate half-buried detection accuracy above 35%.",
                "2. Live Tractor Camera Stream: Build GStreamer / Video4Linux2 real-time camera ingestion pipeline with GPS geo-tagging.",
                "3. Hardware Accelerator Offload: Deploy on Raspberry Pi AI Kit (Hailo-8L NPU, 13 TOPS) for sub-10 ms inference.",
                "4. Closed-Loop Actuator Integration: Interface with hydraulic rock-picker control PLC via CAN bus."
            ]
        },
        "In conclusion, TFLite INT8 is our definitive edge deployment choice. Future work will integrate 3D depth sensing and NPU hardware."
    )

    # Save presentation
    output_path = 'Rock_Detection_Presentation.pptx'
    prs.save(output_path)
    print(f'Successfully built 25-slide presentation at {output_path}')

if __name__ == '__main__':
    build_deck()
